from __future__ import annotations

import hashlib
import multiprocessing
import sys
import time
import zipfile
from dataclasses import FrozenInstanceError
from pathlib import Path
from types import SimpleNamespace
from unittest.mock import AsyncMock

import pytest

from app.lib.mineru import MinerUError, MinerUResult
from app.services.document_quality import PdfDocumentKind
from app.services.document_parser import DocumentParseError, _stop_process, parse_document
from app.services.document_router import ParserEngine, UnsupportedDocumentTypeError
from app.services.markdown_converter import ConversionError


def _clean_markdown() -> str:
    return "# 合同条款\n\n" + "双方应当按照约定履行各自义务。" * 12


def _clean_pdf_text() -> str:
    return "Contract terms and obligations must be performed in good faith. " * 5


def _mineru(markdown: str | None = None) -> SimpleNamespace:
    client = SimpleNamespace()
    client.parse_pdf_async = AsyncMock(
        return_value=MinerUResult(
            task_id="task-safe-id",
            batch_id="batch-safe-id",
            markdown=markdown or _clean_markdown(),
        )
    )
    return client


class _InlineConnection:
    def __init__(self, mailbox: list[tuple[str, str]]) -> None:
        self.mailbox = mailbox

    def send(self, value: tuple[str, str]) -> None:
        self.mailbox.append(value)

    def poll(self) -> bool:
        return bool(self.mailbox)

    def recv(self) -> tuple[str, str]:
        return self.mailbox.pop(0)

    def close(self) -> None:
        pass


class _InlineProcess:
    def __init__(self, target, args) -> None:
        self.target = target
        self.args = args

    def start(self) -> None:
        self.target(*self.args)

    def is_alive(self) -> bool:
        return False

    def terminate(self) -> None:
        pass

    def kill(self) -> None:
        pass

    def join(self, timeout: float | None = None) -> None:
        pass


class _InlineProcessContext:
    def Pipe(self, duplex: bool = False):
        mailbox: list[tuple[str, str]] = []
        return _InlineConnection(mailbox), _InlineConnection(mailbox)

    def Process(self, *, target, args, daemon: bool):
        assert target.__module__ == "app.services.document_parser"
        return _InlineProcess(target, args)


def _write_pdf(path: Path, page_texts: list[str]) -> None:
    import pymupdf

    document = pymupdf.open()
    for text in page_texts:
        page = document.new_page()
        if text:
            page.insert_textbox(page.rect, text, fontsize=10)
    document.save(path)
    document.close()


def _assert_mineru_received_private_snapshot(mineru: SimpleNamespace, source: Path) -> None:
    snapshot = Path(mineru.parse_pdf_async.await_args.args[0])
    assert snapshot != source
    assert not snapshot.exists()


@pytest.mark.asyncio
@pytest.mark.parametrize("suffix", [".txt", ".md"])
async def test_text_and_markdown_are_read_as_utf8_and_normalized(
    tmp_path: Path,
    suffix: str,
) -> None:
    path = tmp_path / f"document{suffix}"
    path.write_bytes(("标题  \r\n\r\n" + "有效正文。" * 20 + "\r\n").encode())

    result = await parse_document(path, suffix, mineru_client=_mineru())

    assert result.markdown == "标题\n\n" + "有效正文。" * 20 + "\n"
    assert result.parser_engine is ParserEngine.DIRECT_TEXT
    assert len(result.content_hash) == 64
    assert len(result.markdown_hash) == 64
    assert result.quality.is_acceptable is True
    assert result.mineru_task_id is None
    with pytest.raises(FrozenInstanceError):
        result.markdown = "changed"  # type: ignore[misc]


@pytest.mark.asyncio
async def test_word_falls_back_to_native_async_mineru_after_quality_gate(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    path = tmp_path / "document.docx"
    path.write_bytes(b"validated word input")
    def low_quality_conversion(_path: str | Path) -> str:
        return "too short"

    monkeypatch.setattr("app.services.document_parser.convert_to_markdown", low_quality_conversion)
    mineru = _mineru()
    result = await parse_document(
        path,
        ".docx",
        mineru_client=mineru,
        process_context=_InlineProcessContext(),
    )

    _assert_mineru_received_private_snapshot(mineru, path)
    assert result.parser_engine is ParserEngine.MINERU
    assert result.mineru_task_id == "task-safe-id"
    assert not hasattr(result, "url")


@pytest.mark.asyncio
@pytest.mark.parametrize("kind", [PdfDocumentKind.SCANNED, PdfDocumentKind.MIXED])
async def test_scanned_and_mixed_pdf_go_directly_to_mineru(
    tmp_path: Path,
    kind: PdfDocumentKind,
) -> None:
    path = tmp_path / "document.pdf"
    _write_pdf(path, [""] if kind is PdfDocumentKind.SCANNED else [_clean_pdf_text(), ""])
    mineru = _mineru()

    result = await parse_document(
        path,
        ".pdf",
        mineru_client=mineru,
        process_context=_InlineProcessContext(),
    )

    assert result.parser_engine is ParserEngine.MINERU
    _assert_mineru_received_private_snapshot(mineru, path)


@pytest.mark.asyncio
async def test_text_pdf_uses_markitdown_then_mineru_when_second_gate_fails(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    path = tmp_path / "document.pdf"
    _write_pdf(path, [_clean_pdf_text()])
    monkeypatch.setattr("app.services.document_parser.convert_to_markdown", lambda _path: "too short")
    mineru = _mineru()

    result = await parse_document(
        path,
        ".pdf",
        mineru_client=mineru,
        process_context=_InlineProcessContext(),
    )

    assert result.parser_engine is ParserEngine.MINERU
    _assert_mineru_received_private_snapshot(mineru, path)


@pytest.mark.asyncio
@pytest.mark.parametrize("suffix", [".pptx", ".xlsx"])
async def test_presentation_and_spreadsheet_failure_has_stable_business_code(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
    suffix: str,
) -> None:
    path = tmp_path / f"document{suffix}"
    path.write_bytes(b"validated office input")

    def fail_conversion(_path: str | Path) -> str:
        raise ConversionError("internal conversion details")

    monkeypatch.setattr("app.services.document_parser.convert_to_markdown", fail_conversion)

    with pytest.raises(DocumentParseError) as caught:
        await parse_document(
            path,
            suffix,
            mineru_client=_mineru(),
            process_context=_InlineProcessContext(),
        )

    assert caught.value.code == "convert_to_pdf_required"
    assert "PDF" in str(caught.value)
    assert "internal conversion details" not in str(caught.value)


@pytest.mark.asyncio
async def test_images_are_rejected_without_mineru(tmp_path: Path) -> None:
    path = tmp_path / "image.png"
    path.write_bytes(b"validated image input")
    mineru = _mineru()

    with pytest.raises(UnsupportedDocumentTypeError):
        await parse_document(path, ".png", mineru_client=mineru)

    mineru.parse_pdf_async.assert_not_awaited()


@pytest.mark.asyncio
async def test_malformed_suffix_is_rejected_before_sensitive_snapshot_creation(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    path = tmp_path / "document.txt"
    path.write_text(_clean_markdown(), encoding="utf-8")
    monkeypatch.setattr(
        "app.services.document_parser.create_private_temp_file",
        lambda **_kwargs: (_ for _ in ()).throw(AssertionError("snapshot must not be created")),
    )

    with pytest.raises(UnsupportedDocumentTypeError):
        await parse_document(path, ".txt/../../private")


@pytest.mark.asyncio
async def test_blocking_conversion_is_bounded_by_total_deadline(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    path = tmp_path / "document.docx"
    path.write_bytes(b"validated word input")

    def slow_conversion(_path: str | Path) -> str:
        time.sleep(0.05)
        return _clean_markdown()

    monkeypatch.setattr("app.services.document_parser.convert_to_markdown", slow_conversion)

    with pytest.raises(DocumentParseError) as caught:
        await parse_document(
            path,
            ".docx",
            mineru_client=_mineru(),
            deadline_seconds=0.001,
        )

    assert caught.value.code == "parse_timeout"


@pytest.mark.asyncio
@pytest.mark.parametrize("suffix", [".txt", ".md"])
async def test_zero_deadline_rejects_snapshot_and_direct_text_read(
    tmp_path: Path,
    suffix: str,
) -> None:
    path = tmp_path / f"private{suffix}"
    path.write_text(_clean_markdown(), encoding="utf-8")

    with pytest.raises(DocumentParseError) as caught:
        await parse_document(path, suffix, deadline_seconds=0)

    assert caught.value.code == "parse_timeout"
    assert caught.value.__cause__ is None


@pytest.mark.asyncio
async def test_real_minimal_text_pdf_is_classified_and_parsed_locally(tmp_path: Path) -> None:
    path = tmp_path / "document.pdf"
    _write_pdf(path, [_clean_pdf_text()])

    result = await parse_document(path, ".pdf", mineru_client=_mineru())

    assert result.parser_engine is ParserEngine.MARKITDOWN
    assert result.quality.is_acceptable is True


@pytest.mark.asyncio
@pytest.mark.parametrize("suffix", [".txt", ".md"])
async def test_low_quality_direct_text_has_stable_terminal_error(
    tmp_path: Path,
    suffix: str,
) -> None:
    path = tmp_path / f"short{suffix}"
    path.write_text("短文本", encoding="utf-8")

    with pytest.raises(DocumentParseError) as caught:
        await parse_document(path, suffix, process_context=_InlineProcessContext())

    assert caught.value.code == "low_quality"
    assert caught.value.__cause__ is None


@pytest.mark.asyncio
@pytest.mark.parametrize("suffix", [".pptx", ".xlsx"])
async def test_low_quality_markitdown_only_result_requires_pdf(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
    suffix: str,
) -> None:
    path = tmp_path / f"short{suffix}"
    path.write_bytes(b"office")
    monkeypatch.setattr("app.services.document_parser.convert_to_markdown", lambda _path: "短文本")

    with pytest.raises(DocumentParseError) as caught:
        await parse_document(path, suffix, process_context=_InlineProcessContext())

    assert caught.value.code == "convert_to_pdf_required"
    assert caught.value.__cause__ is None


@pytest.mark.asyncio
async def test_low_quality_mineru_result_has_stable_terminal_error(tmp_path: Path) -> None:
    path = tmp_path / "scanned.pdf"
    _write_pdf(path, [""])
    mineru = _mineru("短文本")

    with pytest.raises(DocumentParseError) as caught:
        await parse_document(
            path,
            ".pdf",
            mineru_client=mineru,
            process_context=_InlineProcessContext(),
        )

    assert caught.value.code == "mineru_low_quality"
    assert caught.value.__cause__ is None


@pytest.mark.asyncio
async def test_invalid_utf8_has_sanitized_stable_error(tmp_path: Path) -> None:
    path = tmp_path / "secret-name.txt"
    path.write_bytes(b"\xff\xfe")

    with pytest.raises(DocumentParseError) as caught:
        await parse_document(path, ".txt")

    assert caught.value.code == "invalid_utf8"
    assert str(path) not in str(caught.value)
    assert caught.value.__cause__ is None


@pytest.mark.asyncio
async def test_missing_file_has_sanitized_stable_error(tmp_path: Path) -> None:
    path = tmp_path / "private-name.docx"

    with pytest.raises(DocumentParseError) as caught:
        await parse_document(path, ".docx")

    assert caught.value.code == "file_read"
    assert str(path) not in str(caught.value)
    assert caught.value.__cause__ is None


@pytest.mark.asyncio
async def test_snapshot_file_is_cleaned_when_source_open_fails(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    source = tmp_path / "missing.docx"
    created_path: Path | None = None
    real_mkstemp = __import__("tempfile").mkstemp

    def recording_mkstemp(*args, **kwargs):
        nonlocal created_path
        descriptor, path = real_mkstemp(*args, **kwargs)
        created_path = Path(path)
        return descriptor, path

    monkeypatch.setattr("app.lib.private_temp.tempfile.mkstemp", recording_mkstemp)

    with pytest.raises(DocumentParseError):
        await parse_document(source, ".docx")

    assert created_path is not None
    assert not created_path.exists()


@pytest.mark.asyncio
async def test_unexpected_conversion_error_is_sanitized(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    path = tmp_path / "private.pptx"
    path.write_bytes(b"word")
    monkeypatch.setattr(
        "app.services.document_parser.convert_to_markdown",
        lambda _path: (_ for _ in ()).throw(RuntimeError("sensitive converter detail")),
    )

    with pytest.raises(DocumentParseError) as caught:
        await parse_document(path, ".pptx", process_context=_InlineProcessContext())

    assert caught.value.code == "convert_to_pdf_required"
    assert "sensitive" not in str(caught.value)
    assert caught.value.__cause__ is None


@pytest.mark.asyncio
async def test_mineru_error_is_sanitized(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    path = tmp_path / "scan.pdf"
    _write_pdf(path, [""])
    mineru = _mineru()
    mineru.parse_pdf_async.side_effect = MinerUError("upload", "sensitive upstream detail")

    with pytest.raises(DocumentParseError) as caught:
        await parse_document(
            path,
            ".pdf",
            mineru_client=mineru,
            process_context=_InlineProcessContext(),
        )

    assert caught.value.code == "mineru_failed"
    assert "sensitive" not in str(caught.value)
    assert caught.value.__cause__ is None


@pytest.mark.asyncio
async def test_streaming_snapshot_has_exact_hash_and_is_the_only_parsed_file(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    source = tmp_path / "source.docx"
    original = _clean_markdown().encode("utf-8")
    source.write_bytes(original)
    parsed_paths: list[Path] = []

    def replace_source_then_convert(snapshot: str | Path) -> str:
        parsed_paths.append(Path(snapshot))
        source.write_bytes(b"replacement")
        return Path(snapshot).read_text(encoding="utf-8")

    monkeypatch.setattr("app.services.document_parser.convert_to_markdown", replace_source_then_convert)

    result = await parse_document(
        source,
        ".docx",
        mineru_client=_mineru(),
        process_context=_InlineProcessContext(),
    )

    assert result.content_hash == hashlib.sha256(original).hexdigest()
    assert result.markdown == _clean_markdown() + "\n"
    assert parsed_paths[0] != source
    assert not parsed_paths[0].exists()


@pytest.mark.asyncio
async def test_configured_parse_size_limit_is_enforced_while_copying(
    tmp_path: Path,
) -> None:
    path = tmp_path / "large.txt"
    path.write_bytes(b"x" * 11)

    with pytest.raises(DocumentParseError) as caught:
        await parse_document(path, ".txt", max_parse_bytes=10)

    assert caught.value.code == "file_too_large"
    assert caught.value.__cause__ is None


@pytest.mark.asyncio
async def test_doc_extension_can_succeed_with_local_conversion(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    path = tmp_path / "legacy.doc"
    path.write_bytes(b"legacy word")
    mineru = _mineru()
    monkeypatch.setattr("app.services.document_parser.convert_to_markdown", lambda _path: _clean_markdown())

    result = await parse_document(
        path,
        ".doc",
        mineru_client=mineru,
        process_context=_InlineProcessContext(),
    )

    assert result.parser_engine is ParserEngine.MARKITDOWN
    assert result.quality.is_acceptable is True
    mineru.parse_pdf_async.assert_not_awaited()


@pytest.mark.asyncio
async def test_real_minimal_docx_succeeds_locally(tmp_path: Path) -> None:
    path = tmp_path / "word.docx"
    body = "Contract obligations must be performed in good faith. " * 5
    document_xml = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<w:document xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main">'
        f"<w:body><w:p><w:r><w:t>{body}</w:t></w:r></w:p></w:body></w:document>"
    )
    with zipfile.ZipFile(path, "w") as archive:
        archive.writestr("word/document.xml", document_xml)

    result = await parse_document(path, ".docx", mineru_client=_mineru())

    assert result.parser_engine is ParserEngine.MARKITDOWN
    assert body.strip() in result.markdown


@pytest.mark.asyncio
@pytest.mark.parametrize("suffix", [".pptx", ".xlsx"])
async def test_real_minimal_presentation_and_spreadsheet_succeed_locally(
    tmp_path: Path,
    suffix: str,
) -> None:
    path = tmp_path / f"office{suffix}"
    text = "Contract obligations and evidence requirements. " * 8
    if suffix == ".pptx":
        from pptx import Presentation

        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = "Contract"
        slide.placeholders[1].text = text
        presentation.save(path)
    else:
        from openpyxl import Workbook

        workbook = Workbook()
        workbook.active["A1"] = text
        workbook.save(path)

    result = await parse_document(path, suffix, mineru_client=_mineru())

    assert result.parser_engine is ParserEngine.MARKITDOWN
    assert result.quality.is_acceptable is True


@pytest.mark.asyncio
async def test_invalid_pdf_has_sanitized_stable_error(tmp_path: Path) -> None:
    path = tmp_path / "private-invalid.pdf"
    path.write_bytes(b"not a pdf")

    with pytest.raises(DocumentParseError) as caught:
        await parse_document(path, ".pdf", process_context=_InlineProcessContext())

    assert caught.value.code == "pdf_read"
    assert str(path) not in str(caught.value)
    assert caught.value.__cause__ is None


@pytest.mark.asyncio
async def test_deadline_terminates_and_joins_process_without_late_side_effect(tmp_path: Path) -> None:
    path = tmp_path / "word.docx"
    path.write_bytes(b"word")
    late_side_effect = tmp_path / "late.txt"

    class HangingProcess:
        alive = False
        terminated = False
        joined = False

        def start(self) -> None:
            self.alive = True

        def is_alive(self) -> bool:
            return self.alive

        def terminate(self) -> None:
            self.terminated = True
            self.alive = False

        def kill(self) -> None:
            self.alive = False

        def join(self, timeout: float | None = None) -> None:
            self.joined = True

    process = HangingProcess()

    class HangingContext:
        def Pipe(self, duplex: bool = False):
            mailbox: list[tuple[str, str]] = []
            return _InlineConnection(mailbox), _InlineConnection(mailbox)

        def Process(self, **_kwargs):
            return process

    with pytest.raises(DocumentParseError) as caught:
        await parse_document(
            path,
            ".docx",
            deadline_seconds=0.05,
            process_context=HangingContext(),
        )

    await __import__("asyncio").sleep(0.02)
    assert caught.value.code == "parse_timeout"
    assert process.terminated is True
    assert process.joined is True
    assert not late_side_effect.exists()


def _write_late_side_effect(path: str) -> None:
    time.sleep(0.4)
    Path(path).write_text("late", encoding="utf-8")


@pytest.mark.asyncio
@pytest.mark.skipif(sys.platform != "win32", reason="Windows spawn process regression")
async def test_real_windows_spawn_process_has_no_late_side_effect_or_orphan(tmp_path: Path) -> None:
    side_effect = tmp_path / "late-real.txt"
    process = multiprocessing.get_context("spawn").Process(
        target=_write_late_side_effect,
        args=(str(side_effect),),
    )
    process.start()

    await _stop_process(process)
    await __import__("asyncio").sleep(0.6)

    assert process.exitcode is not None
    assert process.is_alive() is False
    assert not side_effect.exists()


@pytest.mark.asyncio
async def test_stop_process_control_method_failures_are_bounded_and_do_not_escape() -> None:
    class BrokenProcess:
        def is_alive(self) -> bool:
            raise OSError("sensitive is_alive failure")

        def terminate(self) -> None:
            raise OSError("sensitive terminate failure")

        def kill(self) -> None:
            raise OSError("sensitive kill failure")

        def join(self, timeout: float | None = None) -> None:
            raise OSError("sensitive join failure")

    started = time.monotonic()
    await _stop_process(BrokenProcess())
    assert time.monotonic() - started < 1


@pytest.mark.asyncio
async def test_stop_process_finishes_kill_and_join_before_propagating_cancellation() -> None:
    class StubbornProcess:
        killed = False
        joined = False

        def is_alive(self) -> bool:
            return not self.killed

        def terminate(self) -> None:
            pass

        def kill(self) -> None:
            self.killed = True

        def join(self, timeout: float | None = None) -> None:
            self.joined = True

    process = StubbornProcess()
    task = __import__("asyncio").create_task(_stop_process(process))
    await __import__("asyncio").sleep(0)
    task.cancel()

    with pytest.raises(__import__("asyncio").CancelledError):
        await task

    assert process.killed is True
    assert process.joined is True


@pytest.mark.asyncio
async def test_isolated_worker_pipe_failure_has_route_specific_sanitized_error(tmp_path: Path) -> None:
    path = tmp_path / "private-slides.pptx"
    path.write_bytes(b"office")

    class PipeFailureContext(_InlineProcessContext):
        pipe_count = 0

        def Pipe(self, duplex: bool = False):
            self.pipe_count += 1
            if self.pipe_count == 2:
                raise OSError("sensitive IPC detail")
            return super().Pipe(duplex=duplex)

    with pytest.raises(DocumentParseError) as caught:
        await parse_document(path, ".pptx", process_context=PipeFailureContext())

    assert caught.value.code == "convert_to_pdf_required"
    assert "sensitive" not in str(caught.value)


@pytest.mark.asyncio
async def test_conversion_worker_exit_has_route_specific_stable_error(tmp_path: Path) -> None:
    path = tmp_path / "slides.pptx"
    path.write_bytes(b"office")

    class ExitedProcess:
        def start(self) -> None:
            pass

        def is_alive(self) -> bool:
            return False

        def terminate(self) -> None:
            pass

        def kill(self) -> None:
            pass

        def join(self, timeout: float | None = None) -> None:
            pass

    class ExitedContext:
        process_count = 0

        def Pipe(self, duplex: bool = False):
            mailbox: list[tuple[str, str]] = []
            return _InlineConnection(mailbox), _InlineConnection(mailbox)

        def Process(self, *, target, args, **_kwargs):
            self.process_count += 1
            if self.process_count == 1:
                return _InlineProcess(target, args)
            return ExitedProcess()

    with pytest.raises(DocumentParseError) as caught:
        await parse_document(path, ".pptx", process_context=ExitedContext())

    assert caught.value.code == "convert_to_pdf_required"
    assert caught.value.__cause__ is None
